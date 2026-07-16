"""
Generate an animated-style NazmOS merchant deck.

Important: python-pptx does not reliably support rich per-object PowerPoint animations
through its public API. This deck uses a professional, robust technique: progressive
build slides with fade transitions. In PowerPoint, each click feels like an animated
reveal while remaining compatible across PowerPoint, Keynote, Google Slides, and PDF export.

Output:
  NazmOS_Merchant_Deck_Animated.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches
from pptx.dml.color import RGBColor

# Reuse the visual system from the static merchant deck.
from create_nazmos_merchant_deck import (  # noqa: F401
    INK, PAPER, PAPER_2, PAPER_3, CASH, CASH_2, LEAK, GOLD, GOLD_2,
    WHATSAPP, WHITE, MUTED, GRAY, DARK_2,
    FONT_DISPLAY, FONT_BODY_FALLBACK, FONT_MONO, FONT_AR,
    SLIDE_W, SLIDE_H,
    set_bg, add_transition, add_brand, add_tag, add_grid, add_text,
    add_rich_line, add_shape, add_line, add_pill, add_bullet_list,
    add_money_card, add_table_row,
)

TOTAL_CONCEPT_SLIDES = 8


def add_build_note(slide, concept: int, stage: int, stages: int, dark: bool = False) -> None:
    color = GRAY if dark else MUTED
    add_text(slide, 10.75, 7.05, 1.85, 0.2, f"CLICK {stage}/{stages}", 7, color, font=FONT_MONO, align=PP_ALIGN.RIGHT)


def add_slide(prs: Presentation, concept_no: int, stage: int, stages: int, dark: bool = False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, INK if dark else PAPER)
    add_transition(slide)
    add_brand(slide, concept_no, TOTAL_CONCEPT_SLIDES, dark=dark)
    add_build_note(slide, concept_no, stage, stages, dark=dark)
    return slide


# -----------------------------------------------------------------------------
# 1. Cover — Build: brand -> headline -> CTA/proof chips
# -----------------------------------------------------------------------------
def cover(prs: Presentation):
    for stage in range(1, 4):
        slide = add_slide(prs, 1, stage, 3)
        add_grid(slide, RGBColor(0xE8, 0xDE, 0xCB))

        if stage >= 1:
            add_text(slide, 0.8, 0.55, 4.2, 0.65, "Nazm·OS", 30, INK, bold=True, font=FONT_DISPLAY)
            add_text(slide, 8.4, 0.62, 4.1, 0.45, "نظم — نظام استرجاع الأرباح", 18, INK, bold=True, font=FONT_AR, align=PP_ALIGN.RIGHT)

        if stage >= 2:
            add_rich_line(
                slide,
                0.8,
                1.78,
                11.2,
                2.9,
                [
                    ("Find the ", INK, True, False),
                    ("cash", CASH, True, True),
                    ("\ntrapped inside\nyour store.", INK, True, False),
                ],
                size=58,
                font=FONT_DISPLAY,
            )

        if stage >= 3:
            add_pill(slide, 9.25, 0.85, 2.0, 0.34, "SAR 43,000", CASH, size=8)
            add_pill(slide, 10.40, 2.18, 1.85, 0.34, "SAR 12,800", GOLD, color=INK, size=8)
            add_pill(slide, 9.55, 3.48, 1.75, 0.34, "SAR 8,400", LEAK, size=8)
            add_shape(slide, 0.8, 5.72, 3.9, 0.055, CASH, radius=False)
            add_text(
                slide,
                0.8,
                5.9,
                8.9,
                0.5,
                "Dead stock. Empty shelves. Margin leaks. NazmOS finds them and helps you recover — on WhatsApp.",
                13.5,
                INK,
                font=FONT_BODY_FALLBACK,
            )
            add_pill(slide, 0.8, 6.55, 2.65, 0.46, "GET A FREE MONEY AUDIT →", CASH, size=8.5)


# -----------------------------------------------------------------------------
# 2. Trust — Build: headline -> what not needed -> what to send
# -----------------------------------------------------------------------------
def trust(prs: Presentation):
    for stage in range(1, 4):
        slide = add_slide(prs, 2, stage, 3)
        add_tag(slide, "You stay in control")
        if stage >= 1:
            add_text(slide, 0.78, 1.02, 9.8, 1.0, "No complicated setup.\nNo risk to your data.", 34, INK, bold=True, font=FONT_DISPLAY)

        if stage >= 2:
            add_shape(slide, 0.78, 2.55, 5.55, 3.75, PAPER_2, radius=True, line=PAPER_3)
            add_text(slide, 1.08, 2.84, 4.7, 0.35, "What you do NOT need", 14, INK, bold=True, font=FONT_DISPLAY)
            add_bullet_list(
                slide,
                1.08,
                3.35,
                4.7,
                [
                    "No POS replacement needed",
                    "No customer names required",
                    "You can remove sensitive details",
                    "Works with your accountant — does not replace them",
                    "No cashiers need to learn new software",
                ],
                size=13,
                gap=0.46,
            )

        if stage >= 3:
            add_shape(slide, 6.95, 2.55, 5.65, 3.75, INK, radius=True)
            add_text(slide, 7.28, 2.85, 4.8, 0.3, "TO START, SEND US", 9, GOLD_2, bold=True, font=FONT_MONO)
            add_text(slide, 7.28, 3.35, 4.8, 1.2, "1. Last 90 days sales export\n2. Current inventory export", 18, WHITE, bold=True, font=FONT_DISPLAY)
            add_text(slide, 7.28, 4.65, 4.8, 0.65, "Excel, CSV, or POS export is fine.\nNo customer names needed.", 12, PAPER_2, font=FONT_BODY_FALLBACK)
            add_shape(slide, 7.28, 5.62, 4.85, 0.045, GOLD_2, radius=False)
            add_text(slide, 7.28, 5.8, 4.8, 0.36, "Get your Money Audit in 48 hours.", 14, CASH_2, bold=True, italic=True, font=FONT_DISPLAY)


# -----------------------------------------------------------------------------
# 3. Money Audit — Build: big number -> breakdown -> next step
# -----------------------------------------------------------------------------
def money_audit(prs: Presentation):
    rows = [
        ("DS", "Cash trapped in dead stock", "SAR 61,000", LEAK),
        ("M", "Margin recovery opportunity", "SAR 18,000", GOLD),
        ("T", "Inventory transfer opportunity", "SAR 9,200", CASH),
        ("S", "Stockouts expected this week", "7 items", LEAK),
    ]
    for stage in range(1, 4):
        slide = add_slide(prs, 3, stage, 3)
        add_tag(slide, "Your first report")
        add_text(slide, 0.78, 1.0, 8.9, 0.75, "The NazmOS Money Audit.", 36, INK, bold=True, font=FONT_DISPLAY)
        add_pill(slide, 0.78, 1.78, 1.65, 0.32, "SAMPLE REPORT", LEAK, size=7.5)

        if stage >= 1:
            add_shape(slide, 0.78, 2.42, 5.55, 4.32, INK, radius=True)
            add_text(slide, 1.12, 2.75, 4.7, 0.3, "MONEY AT RISK", 10, LEAK, bold=True, font=FONT_MONO)
            add_text(slide, 1.1, 3.18, 4.95, 0.9, "SAR 143,000", 42, WHITE, bold=True, font=FONT_DISPLAY)
            add_text(slide, 1.14, 4.13, 4.8, 0.35, "per year · potential recovery", 11, PAPER_2, font=FONT_BODY_FALLBACK)

        if stage >= 2:
            y = 2.75
            for row in rows:
                add_table_row(slide, 7.02, y, *row)
                add_line(slide, 7.0, y + 0.62, 12.42, y + 0.62, PAPER_3, 0.7)
                y += 0.85

        if stage >= 3:
            add_shape(slide, 1.14, 4.92, 4.75, 0.05, GOLD_2, radius=False)
            add_text(slide, 1.14, 5.18, 4.7, 0.65, "This is what we identify as trapped, leaking, or at risk inside your store.", 12.5, PAPER_2, italic=True, font=FONT_DISPLAY)
            add_shape(slide, 7.02, 6.04, 5.4, 0.6, PAPER_2, radius=True, line=PAPER_3)
            add_text(slide, 7.25, 6.18, 2.45, 0.25, "Recommended next step", 9, MUTED, font=FONT_MONO, bold=True)
            add_text(slide, 9.55, 6.15, 2.65, 0.28, "30-Day Pilot · SAR 3,000", 10.5, CASH, font=FONT_MONO, bold=True, align=PP_ALIGN.RIGHT)


# -----------------------------------------------------------------------------
# 4. What NazmOS watches — Build cards one-by-one
# -----------------------------------------------------------------------------
def watch(prs: Presentation):
    quads = [
        (0.78, 2.22, "Cash Leakage", "Money stuck in slow-moving or dead stock that ties up your cash flow.", LEAK),
        (6.85, 2.22, "Stockout Risk", "Fast-moving items about to run out, costing you sales and customers.", GOLD),
        (0.78, 4.68, "Margin Drops", "Supplier costs increasing quietly, eating into your daily profits.", CASH),
        (6.85, 4.68, "Branch Imbalance", "Overstock in one branch while another runs empty. We suggest transfers.", INK),
    ]
    for stage in range(1, 5):
        slide = add_slide(prs, 4, stage, 4)
        add_tag(slide, "What we watch")
        add_text(slide, 0.78, 1.0, 11.2, 0.8, "We monitor the four ways stores lose money.", 33, INK, bold=True, font=FONT_DISPLAY)
        for i, (x, y, title, body, accent) in enumerate(quads, start=1):
            if stage >= i:
                add_shape(slide, x, y, 5.7, 1.95, PAPER_2, radius=True, line=PAPER_3)
                add_shape(slide, x, y, 0.08, 1.95, accent, radius=False)
                add_text(slide, x + 0.38, y + 0.26, 5.0, 0.35, title, 21, INK, bold=True, font=FONT_DISPLAY)
                add_text(slide, x + 0.38, y + 0.86, 4.85, 0.75, body, 12.6, INK, font=FONT_BODY_FALLBACK)


# -----------------------------------------------------------------------------
# 5. WhatsApp COO — Build: context -> message -> approve buttons/result
# -----------------------------------------------------------------------------
def whatsapp(prs: Presentation):
    for stage in range(1, 4):
        slide = add_slide(prs, 5, stage, 3, dark=True)
        add_tag(slide, "The owner experience", dark=True)
        add_text(slide, 0.78, 1.0, 7.0, 0.7, "Approve fixes in one tap.", 38, WHITE, bold=True, font=FONT_DISPLAY)

        if stage >= 1:
            add_text(slide, 0.82, 2.25, 4.35, 1.85, "NazmOS doesn't just show charts.\nIt sends specific actions directly to the owner's WhatsApp.", 15, PAPER_2, font=FONT_BODY_FALLBACK)
            add_text(slide, 0.82, 4.35, 4.25, 0.85, "No app to learn. No daily logins. Just simple decisions that protect cash.", 15, GOLD_2, bold=True, font=FONT_DISPLAY)

        # phone shell always appears from stage 2 onward
        if stage >= 2:
            add_shape(slide, 5.55, 1.35, 3.65, 5.72, RGBColor(0x05, 0x07, 0x06), radius=True, line=RGBColor(0x32, 0x32, 0x32))
            add_shape(slide, 5.77, 1.62, 3.21, 5.2, RGBColor(0x0B, 0x14, 0x1A), radius=True)
            add_shape(slide, 5.77, 1.62, 3.21, 0.68, DARK_2, radius=True)
            add_shape(slide, 6.0, 1.78, 0.42, 0.42, CASH, shape_type=MSO_AUTO_SHAPE_TYPE.OVAL)
            add_text(slide, 6.0, 1.86, 0.42, 0.2, "N", 10, WHITE, bold=True, font=FONT_DISPLAY, align=PP_ALIGN.CENTER)
            add_text(slide, 6.52, 1.75, 1.55, 0.24, "NazmOS", 10.5, WHITE, bold=True, font=FONT_BODY_FALLBACK)
            add_text(slide, 6.52, 2.0, 1.55, 0.22, "online · Business", 7.5, GRAY, font=FONT_BODY_FALLBACK)
            add_shape(slide, 6.02, 2.65, 2.72, 1.55, RGBColor(0x3D, 0x28, 0x20), radius=True)
            add_text(slide, 6.18, 2.8, 2.38, 0.25, "تنبيه نفاد المخزون", 9.5, LEAK, bold=True, font=FONT_AR, align=PP_ALIGN.RIGHT)
            add_text(slide, 6.18, 3.12, 2.38, 0.82, "حليب المراعي قد ينتهي خلال 1.8 يوم.\nالطلب المقترح: 120 وحدة\nالتكلفة: 840 ر.س", 9.2, WHITE, font=FONT_AR, align=PP_ALIGN.RIGHT)
            add_shape(slide, 6.02, 4.42, 2.72, 1.0, RGBColor(0x1A, 0x3A, 0x26), radius=True)
            add_text(slide, 6.18, 4.55, 2.38, 0.65, "Stockout Warning:\nAlmarai Milk may finish in 1.8 days.\nReorder 120 units.", 8.4, WHITE, font=FONT_BODY_FALLBACK)

        if stage >= 3:
            add_shape(slide, 6.02, 5.72, 1.25, 0.42, WHATSAPP, radius=True)
            add_text(slide, 6.02, 5.81, 1.25, 0.16, "✓ موافقة", 8, INK, bold=True, font=FONT_AR, align=PP_ALIGN.CENTER)
            add_shape(slide, 7.48, 5.72, 1.25, 0.42, DARK_2, radius=True)
            add_text(slide, 7.48, 5.81, 1.25, 0.16, "✕ رفض", 8, WHITE, bold=True, font=FONT_AR, align=PP_ALIGN.CENTER)
            add_shape(slide, 9.8, 2.25, 2.8, 3.35, RGBColor(0x13, 0x20, 0x18), radius=True, line=RGBColor(0x22, 0x44, 0x2C))
            add_text(slide, 10.12, 2.55, 2.1, 0.25, "HOW IT WORKS", 8.5, GOLD_2, bold=True, font=FONT_MONO)
            add_bullet_list(slide, 10.12, 3.05, 2.1, ["Detect risk", "Send recommendation", "Owner approves", "Cash is protected"], 12, PAPER_2, WHATSAPP, 0.48)


# -----------------------------------------------------------------------------
# 6. Weekly Money Report — Build pipeline
# -----------------------------------------------------------------------------
def weekly_report(prs: Presentation):
    for stage in range(1, 4):
        slide = add_slide(prs, 6, stage, 3)
        add_tag(slide, "Your weekly proof")
        add_text(slide, 0.78, 1.0, 10.8, 0.7, "Every Thursday: The Money Report.", 34, INK, bold=True, font=FONT_DISPLAY)

        if stage >= 1:
            add_money_card(slide, 0.78, 2.28, 3.75, 2.12, "Money at Risk", "SAR 25,300", LEAK, "identified opportunity")
        if stage >= 2:
            add_text(slide, 4.62, 3.13, 0.25, 0.35, "→", 18, GOLD, bold=True, font=FONT_BODY_FALLBACK)
            add_money_card(slide, 4.92, 2.28, 3.75, 2.12, "Money Approved", "SAR 9,800", GOLD_2, "owner-approved actions")
        if stage >= 3:
            add_text(slide, 8.75, 3.13, 0.25, 0.35, "→", 18, CASH, bold=True, font=FONT_BODY_FALLBACK)
            add_money_card(slide, 9.05, 2.28, 3.75, 2.12, "Money Recovered", "SAR 6,400", CASH_2, "protected or recovered")
            add_shape(slide, 0.78, 5.05, 11.96, 1.35, PAPER_2, radius=True, line=PAPER_3)
            add_text(slide, 1.1, 5.35, 3.3, 0.22, "STOCKOUTS PREVENTED", 8.5, MUTED, bold=True, font=FONT_MONO)
            add_text(slide, 1.1, 5.68, 2.9, 0.42, "3 items", 22, INK, bold=True, font=FONT_DISPLAY)
            add_line(slide, 4.64, 5.27, 4.64, 6.18, PAPER_3, 1)
            add_text(slide, 5.0, 5.35, 3.4, 0.22, "BUSINESS HEALTH SCORE", 8.5, MUTED, bold=True, font=FONT_MONO)
            add_text(slide, 5.0, 5.68, 2.7, 0.42, "82 / 100", 22, CASH, bold=True, font=FONT_DISPLAY)
            add_line(slide, 8.55, 5.27, 8.55, 6.18, PAPER_3, 1)
            add_text(slide, 8.9, 5.35, 3.4, 0.22, "TOP ACTION", 8.5, MUTED, bold=True, font=FONT_MONO)
            add_text(slide, 8.9, 5.67, 3.45, 0.5, "Transfer milk before Friday", 15, INK, bold=True, font=FONT_DISPLAY)


# -----------------------------------------------------------------------------
# 7. Pricing — Build free -> pilot -> annual
# -----------------------------------------------------------------------------
def pricing(prs: Presentation):
    for stage in range(1, 4):
        slide = add_slide(prs, 7, stage, 3)
        add_tag(slide, "Start small, prove value")
        add_text(slide, 0.78, 1.0, 10.8, 0.72, "Simple pricing after we prove value.", 34, INK, bold=True, font=FONT_DISPLAY)

        if stage >= 1:
            add_shape(slide, 0.78, 2.32, 5.6, 4.45, INK, radius=True)
            add_text(slide, 1.12, 2.65, 4.8, 0.28, "STEP 1: FREE MONEY AUDIT", 9, GOLD_2, bold=True, font=FONT_MONO)
            add_text(slide, 1.12, 3.0, 4.7, 0.62, "SAR 0", 31, WHITE, bold=True, font=FONT_DISPLAY)
        if stage >= 2:
            add_text(slide, 1.12, 3.92, 4.8, 0.28, "STEP 2: 30-DAY RECOVERY PILOT", 9, GOLD_2, bold=True, font=FONT_MONO)
            add_text(slide, 1.12, 4.27, 4.7, 0.62, "SAR 3,000", 31, WHITE, bold=True, font=FONT_DISPLAY)
            add_shape(slide, 1.12, 5.22, 4.75, 0.045, GOLD_2, radius=False)
            add_text(slide, 1.12, 5.48, 4.8, 0.72, "Guarantee: if we do not identify at least SAR 9,000 in recoverable opportunities, we refund the pilot fee.", 11.2, CASH_2, italic=True, font=FONT_DISPLAY)
        if stage >= 3:
            add_shape(slide, 6.92, 2.32, 5.82, 4.45, PAPER_2, radius=True, line=PAPER_3)
            add_text(slide, 7.25, 2.65, 4.9, 0.28, "STEP 3: ANNUAL PLANS", 9, MUTED, bold=True, font=FONT_MONO)
            pricing_rows = [
                ("Small Retail", "1 branch", "SAR 6,900/yr"),
                ("Growing Retail", "2–5 branches", "SAR 18,000/yr"),
                ("Large Chains", "6+ branches", "Custom"),
            ]
            y = 3.22
            for name, sub, price in pricing_rows:
                add_text(slide, 7.25, y, 2.55, 0.3, name, 15.5, INK, bold=True, font=FONT_DISPLAY)
                add_text(slide, 7.25, y + 0.35, 2.55, 0.25, sub, 9.5, MUTED, font=FONT_BODY_FALLBACK)
                add_text(slide, 9.55, y + 0.03, 2.75, 0.32, price, 13.5, CASH, bold=True, font=FONT_MONO, align=PP_ALIGN.RIGHT)
                add_line(slide, 7.25, y + 0.82, 12.25, y + 0.82, PAPER_3, 0.8)
                y += 1.1


# -----------------------------------------------------------------------------
# 8. FAQ/CTA — Build FAQ -> final CTA
# -----------------------------------------------------------------------------
def faq_cta(prs: Presentation):
    for stage in range(1, 3):
        slide = add_slide(prs, 8, stage, 2)
        add_tag(slide, "FAQ & next steps")
        add_text(slide, 0.78, 1.0, 10.5, 0.75, "Ready to find your lost cash?", 35, INK, bold=True, font=FONT_DISPLAY)
        if stage >= 1:
            faqs = [
                ("Do I replace my POS?", "No. We work with your exports."),
                ("Do cashiers use it?", "No. Owners/managers get the reports."),
                ("Is data safe?", "Yes. No customer names needed."),
                ("Live regulated services?", "Outside the core Retail Recovery product for now."),
            ]
            y = 2.28
            for q, a in faqs:
                add_text(slide, 0.82, y, 4.95, 0.28, q, 14.5, INK, bold=True, font=FONT_DISPLAY)
                add_text(slide, 0.82, y + 0.34, 4.95, 0.3, a, 11.5, MUTED, font=FONT_BODY_FALLBACK)
                y += 0.95
        if stage >= 2:
            add_shape(slide, 6.45, 2.25, 6.05, 4.25, INK, radius=True)
            add_text(slide, 6.82, 2.58, 5.2, 0.28, "TO START TODAY", 9, GOLD_2, bold=True, font=FONT_MONO)
            add_bullet_list(
                slide,
                6.82,
                3.1,
                5.05,
                ["Send your last 90 days sales file", "Send your current inventory file", "Get your Free Money Audit in 48 hours"],
                size=13,
                color=WHITE,
                bullet_color=GOLD_2,
                gap=0.5,
            )
            add_text(slide, 6.82, 4.95, 2.6, 0.25, "WHATSAPP US", 8.5, PAPER_2, font=FONT_MONO, bold=True)
            add_text(slide, 6.82, 5.32, 3.0, 0.35, "+966 XX XXX XXXX", 18, CASH_2, bold=True, font=FONT_DISPLAY)
            add_text(slide, 6.82, 6.05, 5.2, 0.35, "نظم — نظام استرجاع أرباح المتجر", 13.5, PAPER_2, font=FONT_AR, align=PP_ALIGN.RIGHT)


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    cover(prs)
    trust(prs)
    money_audit(prs)
    watch(prs)
    whatsapp(prs)
    weekly_report(prs)
    pricing(prs)
    faq_cta(prs)
    return prs


if __name__ == "__main__":
    out = Path("NazmOS_Merchant_Deck_Animated.pptx")
    deck = build_deck()
    deck.save(out)
    print(f"Successfully generated {out.resolve()}")
