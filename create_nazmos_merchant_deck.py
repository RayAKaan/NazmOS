"""
Generate the NazmOS customer-facing merchant deck.

Output:
  NazmOS_Merchant_Deck.pptx

Install if needed:
  pip install python-pptx lxml

Positioning:
  NazmOS is a Retail Recovery System: free money audit -> paid recovery pilot -> annual plan.
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# -----------------------------------------------------------------------------
# Brand palette
# -----------------------------------------------------------------------------
INK = RGBColor(0x0A, 0x0E, 0x0C)
PAPER = RGBColor(0xF4, 0xEF, 0xE6)
PAPER_2 = RGBColor(0xEC, 0xE5, 0xD6)
PAPER_3 = RGBColor(0xDF, 0xD5, 0xC2)
CASH = RGBColor(0x0B, 0x6B, 0x3A)
CASH_2 = RGBColor(0x13, 0xA0, 0x5A)
LEAK = RGBColor(0xC8, 0x41, 0x2A)
GOLD = RGBColor(0xB8, 0x86, 0x2E)
GOLD_2 = RGBColor(0xE0, 0xB3, 0x4A)
WHATSAPP = RGBColor(0x25, 0xD3, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x68, 0x69, 0x61)
GRAY = RGBColor(0x86, 0x96, 0xA0)
DARK = RGBColor(0x1A, 0x17, 0x14)
DARK_2 = RGBColor(0x1F, 0x2C, 0x33)

FONT_DISPLAY = "Georgia"
FONT_BODY = "Aptos"
FONT_BODY_FALLBACK = "Arial"
FONT_MONO = "Consolas"
FONT_AR = "Arial"

SLIDE_W = 13.333
SLIDE_H = 7.5


# -----------------------------------------------------------------------------
# Helpers
# -----------------------------------------------------------------------------
def rgb_tuple(color: RGBColor) -> tuple[int, int, int]:
    return (color[0], color[1], color[2])


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_transition(slide) -> None:
    """Add a simple fade transition. Safe to ignore if a viewer does not support it."""
    xml = """
    <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med">
      <p:fade/>
    </p:transition>
    """
    try:
        slide.element.append(etree.fromstring(xml))
    except Exception:
        pass


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int,
    color: RGBColor = INK,
    *,
    bold: bool = False,
    italic: bool = False,
    font: str = FONT_BODY,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
    return box


def add_rich_line(slide, x, y, w, h, parts, size=32, font=FONT_DISPLAY, align=PP_ALIGN.LEFT):
    """parts = [(text, color, bold, italic), ...]"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for text, color, bold, italic in parts:
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
    return box


def add_shape(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: RGBColor,
    *,
    line: RGBColor | None = None,
    radius: bool = True,
    shape_type=None,
):
    if shape_type is None:
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_line(slide, x1, y1, x2, y2, color=INK, width=1.2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_pill(slide, x, y, w, h, text, fill, color=WHITE, size=10, font=FONT_MONO, bold=True):
    add_shape(slide, x, y, w, h, fill, radius=True)
    return add_text(slide, x, y + 0.045, w, h, text, size, color, bold=bold, font=font, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_brand(slide, n: int, total: int = 8, dark: bool = False) -> None:
    color = PAPER if dark else INK
    muted = GRAY if dark else MUTED
    add_text(slide, 0.55, 0.22, 2.4, 0.25, "● NAZMOS", 8, color, font=FONT_MONO, bold=True)
    add_text(slide, 11.65, 0.22, 1.1, 0.25, f"{n:02d} / {total:02d}", 8, muted, font=FONT_MONO, align=PP_ALIGN.RIGHT)


def add_tag(slide, text, dark=False):
    color = GRAY if dark else MUTED
    add_line(slide, 0.72, 0.73, 1.02, 0.73, color, 1.1)
    add_text(slide, 1.12, 0.59, 5, 0.28, text.upper(), 8.5, color, font=FONT_MONO, bold=True)


def add_grid(slide, color=PAPER_2) -> None:
    """Very light grid for depth."""
    for x in [i * 1.0 for i in range(1, 13)]:
        add_line(slide, x, 0.0, x, SLIDE_H, color, 0.25)
    for y in [i * 1.0 for i in range(1, 8)]:
        add_line(slide, 0.0, y, SLIDE_W, y, color, 0.25)


def add_bullet_list(slide, x, y, w, items, size=14, color=INK, bullet_color=CASH, gap=0.38):
    cur_y = y
    for item in items:
        add_text(slide, x, cur_y, 0.25, 0.28, "→", size, bullet_color, bold=True, font=FONT_BODY_FALLBACK)
        add_text(slide, x + 0.32, cur_y, w - 0.32, 0.32, item, size, color, font=FONT_BODY_FALLBACK)
        cur_y += gap


def add_money_card(slide, x, y, w, h, label, amount, accent, sub=None):
    add_shape(slide, x, y, w, h, INK, radius=True)
    add_text(slide, x + 0.25, y + 0.22, w - 0.5, 0.28, label.upper(), 8.5, accent, bold=True, font=FONT_MONO)
    add_text(slide, x + 0.25, y + 0.62, w - 0.5, 0.72, amount, 28, WHITE, bold=True, font=FONT_DISPLAY)
    if sub:
        add_text(slide, x + 0.25, y + 1.38, w - 0.5, 0.32, sub, 9.5, PAPER_2, font=FONT_BODY_FALLBACK)


def add_table_row(slide, x, y, icon, label, value, accent):
    add_shape(slide, x, y, 0.45, 0.45, PAPER_2, shape_type=MSO_AUTO_SHAPE_TYPE.OVAL, line=PAPER_3)
    add_text(slide, x, y + 0.055, 0.45, 0.28, icon, 9, accent, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.65, y + 0.04, 3.7, 0.35, label, 12.5, INK, font=FONT_BODY_FALLBACK)
    add_text(slide, x + 4.15, y + 0.04, 1.75, 0.35, value, 12.5, INK, bold=True, font=FONT_MONO, align=PP_ALIGN.RIGHT)


# -----------------------------------------------------------------------------
# Deck generation
# -----------------------------------------------------------------------------
def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    # ------------------------------------------------------------------
    # Slide 1 — Cover
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_transition(slide)
    add_grid(slide, RGBColor(0xE8, 0xDE, 0xCB))
    add_brand(slide, 1)

    # Floating sample chips
    add_pill(slide, 9.25, 0.85, 2.0, 0.34, "SAR 43,000", CASH, size=8)
    add_pill(slide, 10.40, 2.18, 1.85, 0.34, "SAR 12,800", GOLD, color=INK, size=8)
    add_pill(slide, 9.55, 3.48, 1.75, 0.34, "SAR 8,400", LEAK, size=8)

    add_text(slide, 0.8, 0.55, 4.2, 0.65, "Nazm·OS", 30, INK, bold=True, font=FONT_DISPLAY)
    add_text(slide, 8.4, 0.62, 4.1, 0.45, "نظم — نظام استرجاع الأرباح", 18, INK, bold=True, font=FONT_AR, align=PP_ALIGN.RIGHT)

    add_rich_line(
        slide,
        0.8,
        1.78,
        11.2,
        2.9,
        [
            ("Find the ", INK, True, False),
            ("cash", CASH, True, True),
            ("\ntrapped inside\nyour store.", INK, True, False),
        ],
        size=58,
        font=FONT_DISPLAY,
    )
    add_shape(slide, 0.8, 5.72, 3.9, 0.055, CASH, radius=False)
    add_text(
        slide,
        0.8,
        5.9,
        8.9,
        0.5,
        "Dead stock. Empty shelves. Margin leaks. NazmOS finds them and helps you recover — on WhatsApp.",
        13.5,
        INK,
        font=FONT_BODY_FALLBACK,
    )
    add_pill(slide, 0.8, 6.55, 2.65, 0.46, "GET A FREE MONEY AUDIT →", CASH, size=8.5)

    # ------------------------------------------------------------------
    # Slide 2 — Trust / What we need
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_transition(slide)
    add_brand(slide, 2)
    add_tag(slide, "You stay in control")
    add_text(slide, 0.78, 1.02, 9.8, 1.0, "No complicated setup.\nNo risk to your data.", 34, INK, bold=True, font=FONT_DISPLAY)

    add_shape(slide, 0.78, 2.55, 5.55, 3.75, PAPER_2, radius=True, line=PAPER_3)
    add_text(slide, 1.08, 2.84, 4.7, 0.35, "What you do NOT need", 14, INK, bold=True, font=FONT_DISPLAY)
    add_bullet_list(
        slide,
        1.08,
        3.35,
        4.7,
        [
            "No POS replacement needed",
            "No customer names required",
            "You can remove sensitive details",
            "Works with your accountant — does not replace them",
            "No cashiers need to learn new software",
        ],
        size=13,
        gap=0.46,
    )

    add_shape(slide, 6.95, 2.55, 5.65, 3.75, INK, radius=True)
    add_text(slide, 7.28, 2.85, 4.8, 0.3, "TO START, SEND US", 9, GOLD_2, bold=True, font=FONT_MONO)
    add_text(slide, 7.28, 3.35, 4.8, 1.2, "1. Last 90 days sales export\n2. Current inventory export", 18, WHITE, bold=True, font=FONT_DISPLAY)
    add_text(slide, 7.28, 4.65, 4.8, 0.65, "Excel, CSV, or POS export is fine.\nNo customer names needed.", 12, PAPER_2, font=FONT_BODY_FALLBACK)
    add_shape(slide, 7.28, 5.62, 4.85, 0.045, GOLD_2, radius=False)
    add_text(slide, 7.28, 5.8, 4.8, 0.36, "Get your Money Audit in 48 hours.", 14, CASH_2, bold=True, italic=True, font=FONT_DISPLAY)

    # ------------------------------------------------------------------
    # Slide 3 — Money Audit
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_transition(slide)
    add_brand(slide, 3)
    add_tag(slide, "Your first report")
    add_text(slide, 0.78, 1.0, 8.9, 0.75, "The NazmOS Money Audit.", 36, INK, bold=True, font=FONT_DISPLAY)
    add_pill(slide, 0.78, 1.78, 1.65, 0.32, "SAMPLE REPORT", LEAK, size=7.5)

    add_shape(slide, 0.78, 2.42, 5.55, 4.32, INK, radius=True)
    add_text(slide, 1.12, 2.75, 4.7, 0.3, "MONEY AT RISK", 10, LEAK, bold=True, font=FONT_MONO)
    add_text(slide, 1.1, 3.18, 4.95, 0.9, "SAR 143,000", 42, WHITE, bold=True, font=FONT_DISPLAY)
    add_text(slide, 1.14, 4.13, 4.8, 0.35, "per year · potential recovery", 11, PAPER_2, font=FONT_BODY_FALLBACK)
    add_shape(slide, 1.14, 4.92, 4.75, 0.05, GOLD_2, radius=False)
    add_text(slide, 1.14, 5.18, 4.7, 0.65, "This is what we identify as trapped, leaking, or at risk inside your store.", 12.5, PAPER_2, italic=True, font=FONT_DISPLAY)

    rows = [
        ("DS", "Cash trapped in dead stock", "SAR 61,000", LEAK),
        ("M", "Margin recovery opportunity", "SAR 18,000", GOLD),
        ("T", "Inventory transfer opportunity", "SAR 9,200", CASH),
        ("S", "Stockouts expected this week", "7 items", LEAK),
    ]
    y = 2.75
    for row in rows:
        add_table_row(slide, 7.02, y, *row)
        add_line(slide, 7.0, y + 0.62, 12.42, y + 0.62, PAPER_3, 0.7)
        y += 0.85

    add_shape(slide, 7.02, 6.04, 5.4, 0.6, PAPER_2, radius=True, line=PAPER_3)
    add_text(slide, 7.25, 6.18, 2.45, 0.25, "Recommended next step", 9, MUTED, font=FONT_MONO, bold=True)
    add_text(slide, 9.55, 6.15, 2.65, 0.28, "30-Day Pilot · SAR 3,000", 10.5, CASH, font=FONT_MONO, bold=True, align=PP_ALIGN.RIGHT)

    # ------------------------------------------------------------------
    # Slide 4 — What NazmOS Watches
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_transition(slide)
    add_brand(slide, 4)
    add_tag(slide, "What we watch")
    add_text(slide, 0.78, 1.0, 11.2, 0.8, "We monitor the four ways stores lose money.", 33, INK, bold=True, font=FONT_DISPLAY)

    quads = [
        (0.78, 2.22, "Cash Leakage", "Money stuck in slow-moving or dead stock that ties up your cash flow.", LEAK),
        (6.85, 2.22, "Stockout Risk", "Fast-moving items about to run out, costing you sales and customers.", GOLD),
        (0.78, 4.68, "Margin Drops", "Supplier costs increasing quietly, eating into your daily profits.", CASH),
        (6.85, 4.68, "Branch Imbalance", "Overstock in one branch while another runs empty. We suggest transfers.", INK),
    ]
    for x, y, title, body, accent in quads:
        add_shape(slide, x, y, 5.7, 1.95, PAPER_2, radius=True, line=PAPER_3)
        add_shape(slide, x, y, 0.08, 1.95, accent, radius=False)
        add_text(slide, x + 0.38, y + 0.26, 5.0, 0.35, title, 21, INK, bold=True, font=FONT_DISPLAY)
        add_text(slide, x + 0.38, y + 0.86, 4.85, 0.75, body, 12.6, INK, font=FONT_BODY_FALLBACK)

    # ------------------------------------------------------------------
    # Slide 5 — WhatsApp COO
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_bg(slide, INK)
    add_transition(slide)
    add_brand(slide, 5, dark=True)
    add_tag(slide, "The owner experience", dark=True)
    add_text(slide, 0.78, 1.0, 7.0, 0.7, "Approve fixes in one tap.", 38, WHITE, bold=True, font=FONT_DISPLAY)
    add_text(slide, 0.82, 2.25, 4.35, 1.85, "NazmOS doesn't just show charts.\nIt sends specific actions directly to the owner's WhatsApp.", 15, PAPER_2, font=FONT_BODY_FALLBACK)
    add_text(slide, 0.82, 4.35, 4.25, 0.85, "No app to learn. No daily logins. Just simple decisions that protect cash.", 15, GOLD_2, bold=True, font=FONT_DISPLAY)

    # Phone shell
    add_shape(slide, 5.55, 1.35, 3.65, 5.72, RGBColor(0x05, 0x07, 0x06), radius=True, line=RGBColor(0x32, 0x32, 0x32))
    add_shape(slide, 5.77, 1.62, 3.21, 5.2, RGBColor(0x0B, 0x14, 0x1A), radius=True)
    add_shape(slide, 5.77, 1.62, 3.21, 0.68, DARK_2, radius=True)
    add_shape(slide, 6.0, 1.78, 0.42, 0.42, CASH, shape_type=MSO_AUTO_SHAPE_TYPE.OVAL)
    add_text(slide, 6.0, 1.86, 0.42, 0.2, "N", 10, WHITE, bold=True, font=FONT_DISPLAY, align=PP_ALIGN.CENTER)
    add_text(slide, 6.52, 1.75, 1.55, 0.24, "NazmOS", 10.5, WHITE, bold=True, font=FONT_BODY_FALLBACK)
    add_text(slide, 6.52, 2.0, 1.55, 0.22, "online · Business", 7.5, GRAY, font=FONT_BODY_FALLBACK)

    add_shape(slide, 6.02, 2.65, 2.72, 1.55, RGBColor(0x3D, 0x28, 0x20), radius=True)
    add_text(slide, 6.18, 2.8, 2.38, 0.25, "تنبيه نفاد المخزون", 9.5, LEAK, bold=True, font=FONT_AR, align=PP_ALIGN.RIGHT)
    add_text(slide, 6.18, 3.12, 2.38, 0.82, "حليب المراعي قد ينتهي خلال 1.8 يوم.\nالطلب المقترح: 120 وحدة\nالتكلفة: 840 ر.س", 9.2, WHITE, font=FONT_AR, align=PP_ALIGN.RIGHT)

    add_shape(slide, 6.02, 4.42, 2.72, 1.0, RGBColor(0x1A, 0x3A, 0x26), radius=True)
    add_text(slide, 6.18, 4.55, 2.38, 0.65, "Stockout Warning:\nAlmarai Milk may finish in 1.8 days.\nReorder 120 units.", 8.4, WHITE, font=FONT_BODY_FALLBACK)

    add_shape(slide, 6.02, 5.72, 1.25, 0.42, WHATSAPP, radius=True)
    add_text(slide, 6.02, 5.81, 1.25, 0.16, "✓ موافقة", 8, INK, bold=True, font=FONT_AR, align=PP_ALIGN.CENTER)
    add_shape(slide, 7.48, 5.72, 1.25, 0.42, DARK_2, radius=True)
    add_text(slide, 7.48, 5.81, 1.25, 0.16, "✕ رفض", 8, WHITE, bold=True, font=FONT_AR, align=PP_ALIGN.CENTER)

    add_shape(slide, 9.8, 2.25, 2.8, 3.35, RGBColor(0x13, 0x20, 0x18), radius=True, line=RGBColor(0x22, 0x44, 0x2C))
    add_text(slide, 10.12, 2.55, 2.1, 0.25, "HOW IT WORKS", 8.5, GOLD_2, bold=True, font=FONT_MONO)
    add_bullet_list(slide, 10.12, 3.05, 2.1, ["Detect risk", "Send recommendation", "Owner approves", "Cash is protected"], 12, PAPER_2, WHATSAPP, 0.48)

    # ------------------------------------------------------------------
    # Slide 6 — Weekly Money Report
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_transition(slide)
    add_brand(slide, 6)
    add_tag(slide, "Your weekly proof")
    add_text(slide, 0.78, 1.0, 10.8, 0.7, "Every Thursday: The Money Report.", 34, INK, bold=True, font=FONT_DISPLAY)

    add_money_card(slide, 0.78, 2.28, 3.75, 2.12, "Money at Risk", "SAR 25,300", LEAK, "identified opportunity")
    add_text(slide, 4.62, 3.13, 0.25, 0.35, "→", 18, GOLD, bold=True, font=FONT_BODY_FALLBACK)
    add_money_card(slide, 4.92, 2.28, 3.75, 2.12, "Money Approved", "SAR 9,800", GOLD_2, "owner-approved actions")
    add_text(slide, 8.75, 3.13, 0.25, 0.35, "→", 18, CASH, bold=True, font=FONT_BODY_FALLBACK)
    add_money_card(slide, 9.05, 2.28, 3.75, 2.12, "Money Recovered", "SAR 6,400", CASH_2, "protected or recovered")

    add_shape(slide, 0.78, 5.05, 11.96, 1.35, PAPER_2, radius=True, line=PAPER_3)
    add_text(slide, 1.1, 5.35, 3.3, 0.22, "STOCKOUTS PREVENTED", 8.5, MUTED, bold=True, font=FONT_MONO)
    add_text(slide, 1.1, 5.68, 2.9, 0.42, "3 items", 22, INK, bold=True, font=FONT_DISPLAY)
    add_line(slide, 4.64, 5.27, 4.64, 6.18, PAPER_3, 1)
    add_text(slide, 5.0, 5.35, 3.4, 0.22, "BUSINESS HEALTH SCORE", 8.5, MUTED, bold=True, font=FONT_MONO)
    add_text(slide, 5.0, 5.68, 2.7, 0.42, "82 / 100", 22, CASH, bold=True, font=FONT_DISPLAY)
    add_line(slide, 8.55, 5.27, 8.55, 6.18, PAPER_3, 1)
    add_text(slide, 8.9, 5.35, 3.4, 0.22, "TOP ACTION", 8.5, MUTED, bold=True, font=FONT_MONO)
    add_text(slide, 8.9, 5.67, 3.45, 0.5, "Transfer milk before Friday", 15, INK, bold=True, font=FONT_DISPLAY)

    # ------------------------------------------------------------------
    # Slide 7 — Pilot + Pricing
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_transition(slide)
    add_brand(slide, 7)
    add_tag(slide, "Start small, prove value")
    add_text(slide, 0.78, 1.0, 10.8, 0.72, "Simple pricing after we prove value.", 34, INK, bold=True, font=FONT_DISPLAY)

    add_shape(slide, 0.78, 2.32, 5.6, 4.45, INK, radius=True)
    add_text(slide, 1.12, 2.65, 4.8, 0.28, "STEP 1: FREE MONEY AUDIT", 9, GOLD_2, bold=True, font=FONT_MONO)
    add_text(slide, 1.12, 3.0, 4.7, 0.62, "SAR 0", 31, WHITE, bold=True, font=FONT_DISPLAY)
    add_text(slide, 1.12, 3.92, 4.8, 0.28, "STEP 2: 30-DAY RECOVERY PILOT", 9, GOLD_2, bold=True, font=FONT_MONO)
    add_text(slide, 1.12, 4.27, 4.7, 0.62, "SAR 3,000", 31, WHITE, bold=True, font=FONT_DISPLAY)
    add_shape(slide, 1.12, 5.22, 4.75, 0.045, GOLD_2, radius=False)
    add_text(slide, 1.12, 5.48, 4.8, 0.72, "Guarantee: if we do not identify at least SAR 9,000 in recoverable opportunities, we refund the pilot fee.", 11.2, CASH_2, italic=True, font=FONT_DISPLAY)

    add_shape(slide, 6.92, 2.32, 5.82, 4.45, PAPER_2, radius=True, line=PAPER_3)
    add_text(slide, 7.25, 2.65, 4.9, 0.28, "STEP 3: ANNUAL PLANS", 9, MUTED, bold=True, font=FONT_MONO)
    pricing = [
        ("Small Retail", "1 branch", "SAR 6,900/yr"),
        ("Growing Retail", "2–5 branches", "SAR 18,000/yr"),
        ("Large Chains", "6+ branches", "Custom"),
    ]
    y = 3.22
    for name, sub, price in pricing:
        add_text(slide, 7.25, y, 2.55, 0.3, name, 15.5, INK, bold=True, font=FONT_DISPLAY)
        add_text(slide, 7.25, y + 0.35, 2.55, 0.25, sub, 9.5, MUTED, font=FONT_BODY_FALLBACK)
        add_text(slide, 9.55, y + 0.03, 2.75, 0.32, price, 13.5, CASH, bold=True, font=FONT_MONO, align=PP_ALIGN.RIGHT)
        add_line(slide, 7.25, y + 0.82, 12.25, y + 0.82, PAPER_3, 0.8)
        y += 1.1

    # ------------------------------------------------------------------
    # Slide 8 — FAQ + CTA
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_transition(slide)
    add_brand(slide, 8)
    add_tag(slide, "FAQ & next steps")
    add_text(slide, 0.78, 1.0, 10.5, 0.75, "Ready to find your lost cash?", 35, INK, bold=True, font=FONT_DISPLAY)

    faqs = [
        ("Do I replace my POS?", "No. We work with your exports."),
        ("Do cashiers use it?", "No. Owners/managers get the reports."),
        ("Is data safe?", "Yes. No customer names needed."),
        ("Live regulated services?", "Outside the core Retail Recovery product for now."),
    ]
    y = 2.28
    for q, a in faqs:
        add_text(slide, 0.82, y, 4.95, 0.28, q, 14.5, INK, bold=True, font=FONT_DISPLAY)
        add_text(slide, 0.82, y + 0.34, 4.95, 0.3, a, 11.5, MUTED, font=FONT_BODY_FALLBACK)
        y += 0.95

    add_shape(slide, 6.45, 2.25, 6.05, 4.25, INK, radius=True)
    add_text(slide, 6.82, 2.58, 5.2, 0.28, "TO START TODAY", 9, GOLD_2, bold=True, font=FONT_MONO)
    add_bullet_list(
        slide,
        6.82,
        3.1,
        5.05,
        [
            "Send your last 90 days sales file",
            "Send your current inventory file",
            "Get your Free Money Audit in 48 hours",
        ],
        size=13,
        color=WHITE,
        bullet_color=GOLD_2,
        gap=0.5,
    )
    add_text(slide, 6.82, 4.95, 2.6, 0.25, "WHATSAPP US", 8.5, PAPER_2, font=FONT_MONO, bold=True)
    add_text(slide, 6.82, 5.32, 3.0, 0.35, "+966 XX XXX XXXX", 18, CASH_2, bold=True, font=FONT_DISPLAY)
    add_text(slide, 6.82, 6.05, 5.2, 0.35, "نظم — نظام استرجاع أرباح المتجر", 13.5, PAPER_2, font=FONT_AR, align=PP_ALIGN.RIGHT)

    return prs


if __name__ == "__main__":
    out = Path("NazmOS_Merchant_Deck.pptx")
    deck = build_deck()
    deck.save(out)
    print(f"Successfully generated {out.resolve()}")
